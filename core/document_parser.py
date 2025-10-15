# core/document_parser.py
import os
import pdfplumber          # For PDF text extraction
import docx               # For Word documents
import pandas as pd       # For CSV/Excel files
from pptx import Presentation  # For PowerPoint files
from PIL import Image     # For image processing
from io import BytesIO    # For handling image bytes
import pytesseract        # For OCR (Optical Character Recognition)
import fitz               # PyMuPDF - for PDF image extraction
import shutil             # For file/folder operations

def process_file(file_path, document_type):
    # Get file extension to determine processing method
    ext = os.path.splitext(file_path)[-1].lower()
    output_folder = "extracted_data"

    # Clean up old extracted data to avoid confusion
    if os.path.exists(output_folder):
        shutil.rmtree(output_folder)  # Delete entire folder and contents
    os.makedirs(output_folder, exist_ok=True)  # Create fresh folder

    # Create path for extracted text file (same name, .txt extension)
    txt_file_path = os.path.join(output_folder, os.path.basename(file_path).replace(ext, ".txt"))
    extracted_images = []  # List to store paths of extracted images

    # Route to appropriate extraction function based on file type
    if ext == ".pdf":
        extracted_images = extract_text_and_images_from_pdf(file_path, txt_file_path, output_folder)
    elif ext in [".pptx"]:
        extracted_images = extract_text_and_images_from_pptx(file_path, txt_file_path, output_folder)
    elif ext == ".csv":
        extracted_images = extract_text_from_csv(file_path, txt_file_path)
    elif ext == ".docx":
        extracted_images = extract_text_and_images_from_docx(file_path, txt_file_path, output_folder)
    elif ext in [".jpg", ".jpeg", ".png"]:
        extracted_images = extract_text_from_image(file_path, txt_file_path)
    else:
        # Unsupported file type
        return None, []

    return txt_file_path, extracted_images

def extract_text_and_images_from_pdf(pdf_path, txt_file_path, output_folder):
    all_text = []           # Store text from all pages
    extracted_images = []   # Store image file paths
    
    # Extract text using pdfplumber (better for text extraction)
    with pdfplumber.open(pdf_path) as pdf:
        for i, page in enumerate(pdf.pages):
            text = page.extract_text()  # Extract text from current page
            if text:
                # Add page marker for better organization
                all_text.append(f"--- Page {i+1} ---\n{text}")

    # Save all extracted text to file
    with open(txt_file_path, "w", encoding="utf-8") as f:
        f.write("\n\n".join(all_text))

    # Extract images using PyMuPDF (better for image extraction)
    try:
        doc = fitz.open(pdf_path)
        for i, page in enumerate(doc):
            # Get all images from current page
            for img_index, img in enumerate(page.get_images(full=True)):
                xref = img[0]  # Image reference number
                base_image = doc.extract_image(xref)  # Extract image data
                image_bytes = base_image["image"]     # Get image bytes
                
                # Convert bytes to PIL Image
                image = Image.open(BytesIO(image_bytes))
                
                # Create descriptive filename
                img_path = os.path.join(output_folder, f"{os.path.basename(pdf_path)}_p{i+1}_img{img_index+1}.png")
                image.save(img_path)  # Save image file
                extracted_images.append(img_path)  # Add to list
    except Exception as e:
        print(f"Image extraction failed: {e}")

    return extracted_images

def extract_text_and_images_from_docx(docx_path, txt_file_path, output_folder):
    # Load Word document
    doc = docx.Document(docx_path)
    
    # Extract text from all paragraphs (skip empty ones)
    all_text = [para.text for para in doc.paragraphs if para.text.strip()]

    # Save extracted text to file
    with open(txt_file_path, "w", encoding="utf-8") as f:
        f.write("\n".join(all_text))

    # TODO: Add image extraction from DOCX files
    return []

def extract_text_from_csv(csv_path, txt_file_path):
    try:
        # Try UTF-8 encoding first (most common)
        df = pd.read_csv(csv_path)
    except UnicodeDecodeError:
        # Fallback to Latin-1 encoding if UTF-8 fails
        df = pd.read_csv(csv_path, encoding='ISO-8859-1')

    # Convert DataFrame to string format and save
    with open(txt_file_path, "w", encoding="utf-8") as txt_file:
        txt_file.write(df.to_string(index=False))  # Don't include row numbers

    print(f"Extracted text from CSV: {txt_file_path}")
    return []  # No images in CSV files

def extract_text_and_images_from_pptx(pptx_path, txt_file_path, output_folder):
    presentation = Presentation(pptx_path)
    full_text = []
    
    # Loop through all slides in presentation
    for i, slide in enumerate(presentation.slides):
        # Loop through all shapes in slide (text boxes, etc.)
        for shape in slide.shapes:
            # Check if shape contains text
            if hasattr(shape, "text"):
                full_text.append(shape.text)

    # Save all extracted text to file
    with open(txt_file_path, "w", encoding="utf-8") as f:
        f.write("\n".join(full_text))

    # TODO: Add image extraction from PPTX files
    return []

def extract_text_from_image(image_path, txt_file_path):
    # Load image and convert to grayscale (better for OCR)
    img = Image.open(image_path).convert("L")
    
    # Use Tesseract OCR to extract text from image
    text = pytesseract.image_to_string(img)
    
    # Save extracted text to file
    with open(txt_file_path, "w", encoding="utf-8") as f:
        f.write(text)
    
    # Return the original image path (since it's the source)
    return [image_path]
